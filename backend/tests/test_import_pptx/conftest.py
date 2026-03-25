"""Shared PPTX fixtures for DoView import tests."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches


def _add_rect(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_rgb: str | None = None,
    text: str = "",
    hyperlink_slide=None,
    bold: bool = False,
    font_size: int | None = None,
) -> object:
    """Add a rectangle shape to a slide with optional fill, text, and hyperlink."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_rgb)
    else:
        shape.fill.background()

    if text:
        tf = shape.text_frame
        tf.text = text
        if bold:
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.bold = True
        if font_size:
            from pptx.util import Pt

            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_size)

    if hyperlink_slide is not None:
        click_action = shape.click_action
        click_action.target_slide = hyperlink_slide

    return shape


def _add_arrow(
    slide,
    left: int,
    top: int,
    width: int = 228600,
    height: int = 228600,
    fill_rgb: str = "C8C8C8",
) -> object:
    """Add a rightArrow shape to a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_rgb)
    return shape


def _add_textbox(slide, left: int, top: int, width: int, height: int, text: str = "") -> object:
    """Add a text box to a slide."""
    from pptx.util import Emu as _Emu

    txbox = slide.shapes.add_textbox(_Emu(left), _Emu(top), _Emu(width), _Emu(height))
    if text:
        txbox.text_frame.text = text
    return txbox


@pytest.fixture()
def minimal_doview_pptx(tmp_path):
    """Create a minimal but valid DoView PPTX with 4 slides.

    Slide 0: Overview (3 overview tiles + 1 final outcomes tile)
    Slide 1: Final Outcomes (2 white rects with grey bars)
    Slide 2: Outcomes Map (4 outcome boxes in 2 columns, 1 arrow)
    Slide 3: Info page (text only, should be skipped)
    """
    prs = Presentation()
    # Standard slide width/height
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # --- Slide 0: Overview ---
    s0 = prs.slides.add_slide(blank_layout)

    # --- Slide 1: Final Outcomes ---
    s1 = prs.slides.add_slide(blank_layout)

    # --- Slide 2: Outcomes Map ---
    s2 = prs.slides.add_slide(blank_layout)

    # --- Slide 3: Info ---
    s3 = prs.slides.add_slide(blank_layout)

    # Now add shapes with hyperlinks (slides must exist first)

    # Slide 0: Overview tiles
    _add_rect(s0, 457200, 274320, 8229600, 548640, text="My DoView Model")  # title
    _add_rect(
        s0, 3582000, 1496520, 1980000, 720000,
        fill_rgb="FFFFFF", text="Final Outcomes", bold=True,
        hyperlink_slide=s1,
    )
    _add_rect(  # grey bar for final outcomes
        s0, 3582000, 1496520, 1980000, 18000, fill_rgb="BEBEBE",
    )
    _add_rect(
        s0, 1242000, 2648520, 1980000, 720000,
        fill_rgb="FFFFBA", text="Topic Alpha",
        hyperlink_slide=s2,
    )
    _add_rect(
        s0, 3582000, 2648520, 1980000, 720000,
        fill_rgb="F9D3D4", text="Topic Beta",
        hyperlink_slide=s2,  # same target for simplicity
    )
    _add_rect(
        s0, 5922000, 2648520, 1980000, 720000,
        fill_rgb="9FE1FF", text="Topic Gamma",
        hyperlink_slide=s2,
    )
    _add_textbox(s0, 209485, 6537960, 8869680, 320040, "Footer text here")

    # Slide 1: Final Outcomes
    _add_rect(s1, 137160, 137160, 1645920, 548640, fill_rgb="E6E6E6",
              text="Back to Overview", hyperlink_slide=s0)
    _add_rect(s1, 457200, 868680, 8229600, 411480,
              fill_rgb="FFFFFF", text="Final Outcomes", bold=True)
    _add_rect(s1, 457200, 868680, 8229600, 18000, fill_rgb="BEBEBE")  # grey bar
    _add_rect(s1, 685800, 1463040, 7772400, 502920,
              fill_rgb="FFFFFF", text="Outcome A achieved")
    _add_rect(s1, 685800, 1463040, 7772400, 18000, fill_rgb="BEBEBE")
    _add_rect(s1, 685800, 2212560, 7772400, 502920,
              fill_rgb="FFFFFF", text="Outcome B achieved")
    _add_rect(s1, 685800, 2212560, 7772400, 18000, fill_rgb="BEBEBE")
    _add_textbox(s1, 274320, 6555960, 8869680, 320040, "Footer disclaimer")

    # Slide 2: Outcomes Map with 2 columns and arrow
    _add_rect(s2, 137160, 137160, 1645920, 548640, fill_rgb="E6E6E6",
              text="Back to Overview", hyperlink_slide=s0)
    _add_rect(s2, 457200, 868680, 8229600, 411480,
              fill_rgb="FFFFBA", text="Topic Alpha")  # page title
    # Column 1
    _add_rect(s2, 685800, 1515110, 1115568, 731520,
              fill_rgb="FFFFBA", text="Step one done")
    _add_rect(s2, 685800, 2429510, 1115568, 731520,
              fill_rgb="FFFFBA", text="Step two done")
    # Arrow between columns
    _add_arrow(s2, 1938528, 2100000)
    # Column 2
    _add_rect(s2, 2350008, 1463040, 1115568, 835660,
              fill_rgb="FFFFBA", text="Result one")
    _add_rect(s2, 2350008, 2481580, 1115568, 731520,
              fill_rgb="FFFFBA", text="Result two")
    _add_textbox(s2, 274320, 6555960, 8869680, 320040, "Footer")

    # Slide 3: Info page (text only)
    _add_textbox(s3, 457200, 274320, 8229600, 548640, "What is a DoView?")
    _add_textbox(s3, 914400, 1097280, 7315200, 5293757,
                 "A DoView is a diagram used to clarify logic...")
    _add_rect(s3, 137160, 137160, 1645920, 548640, fill_rgb="E6E6E6",
              text="Back to Overview", hyperlink_slide=s0)

    path = tmp_path / "test_doview.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture()
def non_doview_pptx(tmp_path):
    """Create a non-DoView PPTX (plain presentation, no overview tiles)."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank_layout)
    _add_textbox(s1, 457200, 274320, 8229600, 548640, "Hello World")
    _add_textbox(s1, 914400, 1097280, 7315200, 2000000, "This is a normal presentation")

    s2 = prs.slides.add_slide(blank_layout)
    _add_textbox(s2, 457200, 274320, 8229600, 548640, "Slide Two")
    _add_rect(s2, 1000000, 1000000, 2000000, 1000000, fill_rgb="FF0000", text="A red box")

    path = tmp_path / "test_not_doview.pptx"
    prs.save(str(path))
    return str(path)
