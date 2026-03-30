"""Tests for file text extraction service (ADR-115)."""

from __future__ import annotations

import csv
import io
import json

import pytest

from app.ai.extract import ExtractionResult, extract_text

# ---------------------------------------------------------------------------
# Plain text / code files
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_extract_plain_text():
    content = b"Hello, this is a plain text file.\nLine two."
    result = await extract_text("readme.txt", content, "text/plain")
    assert isinstance(result, ExtractionResult)
    assert result.extracted_text == "Hello, this is a plain text file.\nLine two."
    assert result.filename == "readme.txt"
    assert result.size_bytes == len(content)
    assert result.truncated is False
    assert result.error is None


@pytest.mark.asyncio
async def test_extract_markdown():
    content = b"# Title\n\nSome **bold** text."
    result = await extract_text("notes.md", content, "text/markdown")
    assert "# Title" in result.extracted_text
    assert result.error is None


@pytest.mark.asyncio
async def test_extract_json():
    data = {"key": "value", "number": 42}
    content = json.dumps(data).encode()
    result = await extract_text("data.json", content, "application/json")
    assert "key" in result.extracted_text
    assert "value" in result.extracted_text
    assert result.error is None


@pytest.mark.asyncio
async def test_extract_csv():
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["Name", "Age", "City"])
    writer.writerow(["Alice", "30", "Wellington"])
    writer.writerow(["Bob", "25", "Auckland"])
    content = buf.getvalue().encode()

    result = await extract_text("people.csv", content, "text/csv")
    assert "Alice" in result.extracted_text
    assert "Wellington" in result.extracted_text
    assert result.error is None


@pytest.mark.asyncio
async def test_extract_xml():
    content = b"<root><item>Hello</item></root>"
    result = await extract_text("data.xml", content, "application/xml")
    assert "<root>" in result.extracted_text
    assert result.error is None


@pytest.mark.asyncio
async def test_extract_code_file():
    content = b"def hello():\n    print('hello')\n"
    result = await extract_text("main.py", content, "text/x-python")
    assert "def hello" in result.extracted_text
    assert result.error is None


@pytest.mark.asyncio
async def test_extract_yaml():
    content = b"name: test\nvalue: 42\n"
    result = await extract_text("config.yaml", content, "application/x-yaml")
    assert "name: test" in result.extracted_text
    assert result.error is None


# ---------------------------------------------------------------------------
# PDF extraction
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_extract_pdf():
    """Test PDF extraction using a minimal valid PDF."""
    # Minimal PDF with text "Hello World"
    # This is the smallest valid PDF with extractable text
    pdf_bytes = (
        b"%PDF-1.0\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Contents 4 0 R"
        b"/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
        b"4 0 obj<</Length 44>>stream\n"
        b"BT /F1 12 Tf 100 700 Td (Hello World) Tj ET\n"
        b"endstream\nendobj\n"
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"0000000266 00000 n \n"
        b"0000000360 00000 n \n"
        b"trailer<</Size 6/Root 1 0 R>>\n"
        b"startxref\n430\n%%EOF"
    )
    result = await extract_text("doc.pdf", pdf_bytes, "application/pdf")
    # pdfplumber should extract "Hello World" from the PDF
    assert "Hello" in result.extracted_text
    assert result.error is None


# ---------------------------------------------------------------------------
# DOCX extraction
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_extract_docx():
    """Test DOCX extraction using python-docx to create a minimal file."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("First paragraph of the document.")
    doc.add_paragraph("Second paragraph with details.")
    buf = io.BytesIO()
    doc.save(buf)
    content = buf.getvalue()

    result = await extract_text("report.docx", content, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    assert "First paragraph" in result.extracted_text
    assert "Second paragraph" in result.extracted_text
    assert result.error is None


# ---------------------------------------------------------------------------
# XLSX extraction
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_extract_xlsx():
    """Test XLSX extraction using openpyxl to create a minimal file."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["Name", "Score"])
    ws.append(["Alice", 95])
    ws.append(["Bob", 87])
    buf = io.BytesIO()
    wb.save(buf)
    content = buf.getvalue()

    result = await extract_text("scores.xlsx", content, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    assert "Alice" in result.extracted_text
    assert "95" in result.extracted_text
    assert result.error is None


# ---------------------------------------------------------------------------
# PPTX extraction
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_extract_pptx():
    """Test PPTX extraction using python-pptx to create a minimal file."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Slide body content here."
    buf = io.BytesIO()
    prs.save(buf)
    content = buf.getvalue()

    result = await extract_text("deck.pptx", content, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    assert "Slide Title" in result.extracted_text
    assert "body content" in result.extracted_text
    assert result.error is None


# ---------------------------------------------------------------------------
# Edge cases
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_extract_empty_file():
    result = await extract_text("empty.txt", b"", "text/plain")
    assert result.extracted_text == ""
    assert result.error is None


@pytest.mark.asyncio
async def test_extract_binary_fallback():
    """Binary file that can't be decoded as UTF-8 should return an error."""
    content = bytes(range(256)) * 10  # Non-UTF-8 binary data
    result = await extract_text("binary.bin", content, "application/octet-stream")
    assert result.error is not None
    assert result.extracted_text == ""


@pytest.mark.asyncio
async def test_extract_truncation():
    """Text exceeding the max char limit should be truncated."""
    long_text = "A" * 200_000
    result = await extract_text("huge.txt", long_text.encode(), "text/plain")
    assert result.truncated is True
    assert len(result.extracted_text) <= 100_001  # 100k + possible partial


@pytest.mark.asyncio
async def test_extract_unknown_extension_text_content():
    """Unknown extension but valid UTF-8 text content should extract successfully."""
    content = b"This is readable text in an unknown format."
    result = await extract_text("data.custom", content, "application/x-unknown")
    assert "readable text" in result.extracted_text
    assert result.error is None


@pytest.mark.asyncio
async def test_extract_size_bytes():
    """size_bytes should reflect the original file size."""
    content = b"Hello world"
    result = await extract_text("test.txt", content, "text/plain")
    assert result.size_bytes == len(content)
