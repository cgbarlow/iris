"""Low-level PPTX parsing for DoView model import."""

from __future__ import annotations

from dataclasses import dataclass


@dataclass
class PptxShape:
    """A single shape extracted from a PPTX slide."""

    slide_index: int
    shape_id: int
    name: str
    text: str
    left: int  # EMU
    top: int  # EMU
    width: int  # EMU
    height: int  # EMU
    fill_color: str | None  # hex e.g. "FFFFBA" or None
    preset_geometry: str | None  # e.g. "rightArrow", "rect"
    hyperlink_slide_index: int | None  # resolved target slide 0-based index
    is_textbox: bool
    is_picture: bool


@dataclass
class PptxSlide:
    """All shapes extracted from a single PPTX slide."""

    index: int  # 0-based
    shapes: list[PptxShape]


def read_pptx(file_path: str) -> list[PptxSlide]:
    """Parse a PPTX file and return slides with their shapes.

    Extracts position, fill colour, preset geometry, and hyperlink
    targets for each shape.  Skips picture elements and connector
    shapes (decorative lines).
    """
    from pptx import Presentation  # noqa: PLC0415
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

    prs = Presentation(file_path)

    # Build slide-object → 0-based index mapping for hyperlink resolution
    slide_index_map: dict[int, int] = {}
    for idx, slide in enumerate(prs.slides):
        slide_index_map[id(slide)] = idx

    slides: list[PptxSlide] = []
    for slide_idx, slide in enumerate(prs.slides):
        shapes: list[PptxShape] = []
        for shape in slide.shapes:
            # Skip pictures and connectors (decorative in DoView PPTX)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                continue
            # Skip group shapes (not used in DoView)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
            # Skip connector shapes (<p:cxnSp>)
            if hasattr(shape, "begin_x"):
                continue

            # Extract text
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()

            # Extract fill colour
            fill_color = _extract_fill_color(shape)

            # Extract preset geometry
            preset_geometry = _extract_preset_geometry(shape)

            # Detect textbox
            is_textbox = shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX

            # Resolve hyperlink to target slide index
            hyperlink_slide_index = _resolve_hyperlink(shape, slide_index_map)

            shapes.append(
                PptxShape(
                    slide_index=slide_idx,
                    shape_id=shape.shape_id,
                    name=shape.name or "",
                    text=text,
                    left=shape.left or 0,
                    top=shape.top or 0,
                    width=shape.width or 0,
                    height=shape.height or 0,
                    fill_color=fill_color,
                    preset_geometry=preset_geometry,
                    hyperlink_slide_index=hyperlink_slide_index,
                    is_textbox=is_textbox,
                    is_picture=False,
                )
            )
        slides.append(PptxSlide(index=slide_idx, shapes=shapes))
    return slides


def _extract_fill_color(shape: object) -> str | None:
    """Extract the solid fill colour as a hex string, or None."""
    try:
        fill = shape.fill  # type: ignore[attr-defined]
        if fill.type is not None:
            from pptx.enum.dml import MSO_THEME_COLOR  # noqa: PLC0415, F401
            from pptx.dml.color import RGBColor  # noqa: PLC0415

            rgb = fill.fore_color.rgb
            if isinstance(rgb, RGBColor):
                return str(rgb)
    except (AttributeError, TypeError, ValueError):
        pass
    return None


def _extract_preset_geometry(shape: object) -> str | None:
    """Extract the preset geometry name (e.g. 'rightArrow', 'rect')."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:  # type: ignore[attr-defined]
            # auto_shape_type gives the MSO_AUTO_SHAPE_TYPE enum
            auto_type = shape.auto_shape_type  # type: ignore[attr-defined]
            if auto_type is not None:
                return str(auto_type).split("(")[0].strip().split(".")[-1].lower()
    except (AttributeError, TypeError, ValueError):
        pass

    # Fallback: read from XML
    try:
        sp_el = shape._element  # type: ignore[attr-defined]  # noqa: SLF001
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        prst_geom = sp_el.find(".//a:prstGeom", ns)
        if prst_geom is not None:
            return prst_geom.get("prst")
    except (AttributeError, TypeError):
        pass
    return None


def _resolve_hyperlink(shape: object, slide_index_map: dict[int, int]) -> int | None:
    """Resolve a shape's click-action hyperlink to a 0-based slide index."""
    try:
        click_action = shape.click_action  # type: ignore[attr-defined]
        if click_action is not None and click_action.target_slide is not None:
            target_slide = click_action.target_slide
            return slide_index_map.get(id(target_slide))
    except (AttributeError, TypeError, ValueError):
        pass
    return None
