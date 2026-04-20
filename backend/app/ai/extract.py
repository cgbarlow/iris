"""Text extraction from uploaded files for AI context (ADR-115).

Stateless extraction — receives file bytes, returns extracted text.
No database access, no file storage.
"""

from __future__ import annotations

import csv
import io
import logging
from dataclasses import dataclass, field

log = logging.getLogger(__name__)

MAX_EXTRACTED_CHARS = 100_000

# File extensions treated as plain text (UTF-8 decode)
_TEXT_EXTENSIONS = frozenset({
    ".txt", ".md", ".markdown", ".json", ".xml", ".yaml", ".yml",
    ".csv", ".tsv", ".log", ".ini", ".toml", ".cfg", ".conf",
    ".py", ".js", ".ts", ".svelte", ".html", ".css", ".scss",
    ".java", ".c", ".cpp", ".h", ".go", ".rs", ".rb", ".sh",
    ".bat", ".ps1", ".sql", ".r", ".swift", ".kt", ".scala",
})


@dataclass
class ExtractionResult:
    """Result of text extraction from an uploaded file."""

    filename: str
    content_type: str
    size_bytes: int
    extracted_text: str = ""
    truncated: bool = False
    error: str | None = None


def _truncate(text: str) -> tuple[str, bool]:
    """Truncate text to MAX_EXTRACTED_CHARS if needed."""
    if len(text) > MAX_EXTRACTED_CHARS:
        return text[:MAX_EXTRACTED_CHARS], True
    return text, False


def _ext(filename: str) -> str:
    """Return lowercase file extension including the dot."""
    dot = filename.rfind(".")
    if dot == -1:
        return ""
    return filename[dot:].lower()


async def extract_text(
    filename: str,
    content: bytes,
    content_type: str,
) -> ExtractionResult:
    """Extract text from file bytes. Returns ExtractionResult."""
    size_bytes = len(content)
    ext = _ext(filename)

    try:
        text = _dispatch_extract(ext, content, content_type)
        text, truncated = _truncate(text)
        return ExtractionResult(
            filename=filename,
            content_type=content_type,
            size_bytes=size_bytes,
            extracted_text=text,
            truncated=truncated,
        )
    except Exception:
        log.exception("Text extraction failed for %s", filename)
        return ExtractionResult(
            filename=filename,
            content_type=content_type,
            size_bytes=size_bytes,
            error=f"Could not extract text from {filename}",
        )


def _dispatch_extract(ext: str, content: bytes, content_type: str) -> str:
    """Route extraction to the appropriate handler based on extension/content-type."""
    if ext == ".pdf" or content_type == "application/pdf":
        return _extract_pdf(content)

    if ext == ".docx" or "wordprocessingml" in content_type:
        return _extract_docx(content)

    if ext == ".xlsx" or "spreadsheetml" in content_type:
        return _extract_xlsx(content)

    if ext == ".pptx" or "presentationml" in content_type:
        return _extract_pptx(content)

    if ext == ".csv" or content_type == "text/csv":
        return _extract_csv(content)

    # Text-like files: known extensions or text/* content types
    if ext in _TEXT_EXTENSIONS or content_type.startswith("text/"):
        return content.decode("utf-8")

    # Unknown type — attempt UTF-8 decode as fallback
    return _extract_text_fallback(content)


def _extract_pdf(content: bytes) -> str:
    """Extract text from PDF using pdfplumber."""
    import pdfplumber  # noqa: PLC0415

    pages: list[str] = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    return "\n\n".join(pages)


def _extract_docx(content: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document  # noqa: PLC0415

    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_xlsx(content: bytes) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook  # noqa: PLC0415

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append("\t".join(cells))
        if rows:
            header = f"Sheet: {sheet.title}" if len(wb.worksheets) > 1 else ""
            if header:
                parts.append(f"{header}\n{chr(10).join(rows)}")
            else:
                parts.append("\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(io.BytesIO(content))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_csv(content: bytes) -> str:
    """Extract text from CSV."""
    text = content.decode("utf-8")
    reader = csv.reader(io.StringIO(text))
    rows = ["\t".join(row) for row in reader]
    return "\n".join(rows)


def _extract_text_fallback(content: bytes) -> str:
    """Attempt UTF-8 decode for unknown file types."""
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError as exc:
        msg = "File contains binary data that cannot be read as text"
        raise ValueError(msg) from exc
